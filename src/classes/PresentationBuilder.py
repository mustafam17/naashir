from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN

from pathlib import Path
import re


class PresentationBuilder:
    def __init__(
        self,
        title: str,
        speaker_name_ar: str,
        speaker_name_en: str,
        text_ar: str,
        text_en: str,
    ):
        self.title = title
        self.speaker_name_ar = speaker_name_ar
        self.speaker_name_en = speaker_name_en
        self.text_ar = text_ar
        self.text_en = text_en

    def style_arabic_text(arabic_text, placeholder):
        """Styles Arabic text inside placeholders, applying specific format only to Qur'anic verses."""
        # Split text into parts: normal and Qur'anic (text between ﴾ and ﴿)
        parts = re.split(r"(﴾.*?﴿)", arabic_text)

        # Clear the current text in the placeholder
        placeholder.text = ""  # This resets the placeholder's text

        # Add runs with specific styles
        for part in parts:
            run = placeholder.text_frame.paragraphs[0].add_run()

            if part.startswith("﴾") and part.endswith("﴿"):
                print("detected" + part)
                # Apply Qur'anic styling: Traditional Arabic, Bold, #7ABC32
                run.text = part
                run.font.name = "Traditional Arabic"
                run.font.bold = True
                run.font.color.rgb = RGBColor(0x7A, 0xBC, 0x32)  # #7ABC32 (green)
            else:
                # For normal text, just set the text and use the default master slide style
                run.text = part

        # Ensure the text is right-aligned (for Arabic)
        # placeholder.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT

    def create_presentation_from_text(
        self,
        presentation_path: Path,
        output_path: Path,
        text_pairs: list[tuple[str, str]],
    ):
        presentation = Presentation(presentation_path)
        slide_layouts = presentation.slide_master.slide_layouts

        for arabic_text, english_text in text_pairs:
            slide = presentation.slides.add_slide(slide_layouts.get_by_name("Slide"))

            speaker_name_arabic_placeholder = slide.placeholders[10]
            speaker_name_english_placeholder = slide.placeholders[11]
            text_arabic_placeholder = slide.placeholders[12]
            text_english_placeholder = slide.placeholders[13]

            speaker_name_arabic_placeholder.text = self.speaker_name_ar
            speaker_name_english_placeholder.text = self.speaker_name_en

            text_arabic_placeholder.text = arabic_text
            # style_arabic_text(arabic_text, text_arabic_placeholder)

            text_english_placeholder.text = english_text

        presentation.save(output_path)
